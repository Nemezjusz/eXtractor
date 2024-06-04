from flask import Flask, render_template, request, redirect, send_file
from openai import OpenAI
from constants import API_KEY
import os
import fitz  # PyMuPDF
import werkzeug
from pptx import Presentation
import re
client = OpenAI(api_key=API_KEY)

app = Flask(__name__)

app.config['UPLOAD_FOLDER'] = 'uploads/'
app.config['ALLOWED_EXTENSIONS'] = {'pdf', 'docx', 'txt'}
app.config['PPTX_FOLDER'] = 'presentations/'

os.makedirs(app.config['PPTX_FOLDER'], exist_ok=True)


def extract_info(transcript):
    completion = client.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system",
             "content": "You are an information extraction assistant, skilled in extracting information"},
            {"role": "user", "content": f"Znajdz informacje w następujacym tekście: {transcript}. "
                                        f"Zwróć poszukiwane informacje w nastepującym formacie "

                                        f"1. Nazwa Klienta: <uzupełnij>"
                                        f"2. Adres Klienta: <uzupełnij>"
                                        f"3. Typ prztargu: <uzupełnij>"
                                        f"4. Koszty przetargu: <uzupełnij>"
                                        f"5. Model płatności: <uzupełnij>"
                                        f"6. Liczba użytkowników: <uzupełnij>"
                                        f"7. Termin Składania Ofert: <uzupełniji>"
                                        f"8. Termin Realizacji Projektu: <uzupełniji>"}
        ]
    )
    return completion.choices[0].message.content


def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']


def read_pdf_text(pdf_path):
    """Reads text from a PDF file."""
    text = ""
    try:
        # Open the PDF file
        pdf_document = fitz.open(pdf_path)

        # Iterate over each page
        for page_num in range(pdf_document.page_count):
            # Get the page
            page = pdf_document[page_num]
            # Extract text from the page
            text += page.get_text()

        pdf_document.close()
    except Exception as e:
        print(f"Error reading PDF file: {e}")

    return text


def format_extracted_text(text):
    """Formats the extracted text so that every numbered item starts from a new line."""
    formatted_text = re.sub(r'(\d+\.\s*)', r'<br>\1', text)
    formatted_text = formatted_text.lstrip('<br>')
    return formatted_text

def format_extracted_text_presentation(text):
    """Formats the extracted text so that every numbered item starts from a new line."""
    formatted_text = re.sub(r'(\d+\.\s*)', r'\n\1', text)
    formatted_text = formatted_text.lstrip('\n')
    return formatted_text

def create_presentation(text, filename):
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    text_chunks = text.split('\n\n')  # Split text into chunks

    for chunk in text_chunks:
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = "Extracted Information"
        content.text = chunk

    pptx_path = os.path.join(app.config['PPTX_FOLDER'], filename)
    prs.save(pptx_path)
    return pptx_path


@app.route('/')
def index():
    return render_template("main.html", transcript="Text will appear here.", pptx_file=None)


@app.route('/upload', methods=['POST'])
def upload_file():
    if 'audiofile' not in request.files:
        return redirect(request.url)

    file = request.files['audiofile']
    if file.filename == '':
        return redirect(request.url)

    if file and allowed_file(file.filename):
        filename = werkzeug.utils.secure_filename(file.filename)
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(filepath)

        if filename.lower().endswith('.pdf'):
            text = read_pdf_text(filepath)
        else:
            with open(filepath, 'r', encoding='utf-8') as f:
                text = f.read()

        extracted_text = extract_info(text)
        formatted_text = format_extracted_text(extracted_text)
        pptx_filename = filename.rsplit('.', 1)[0] + '.pptx'
        formatted_text2 = format_extracted_text_presentation(extracted_text)
        pptx_path = create_presentation(formatted_text2, pptx_filename)
        return render_template("main.html", transcript=formatted_text, pptx_file=pptx_filename)

    return redirect(request.url)


@app.route('/download/<filename>')
def download_file(filename):
    return send_file(os.path.join(app.config['PPTX_FOLDER'], filename), as_attachment=True)


if __name__ == '__main__':
    app.run()
